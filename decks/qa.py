"""Geometry QA on the generated deck, read back from the file itself.

Renders are unavailable in this environment, so instead of eyeballing images
this reads the real shape boxes out of the .pptx and checks the things a
visual pass would catch: content escaping the slide, text boxes that are too
short for the text in them, and cards overlapping each other.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_IN = 914400.0
SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.5          # minimum breathing room from the slide edge
CHAR_W = {            # rough advance width as a fraction of font size
    'Georgia': 0.50, 'Calibri': 0.46, 'Consolas': 0.55,
}

prs = Presentation('ARAH-System-Overview.pptx')
issues = []


def box(sh):
    return (sh.left / EMU_IN, sh.top / EMU_IN,
            sh.width / EMU_IN, sh.height / EMU_IN)


def est_lines(text, width_in, size_pt, face):
    """Estimate wrapped line count for a run of text at a given box width."""
    if not text.strip():
        return 0
    per_char = CHAR_W.get(face, 0.48) * size_pt / 72.0
    if per_char <= 0:
        return 0
    cols = max(1, int(width_in / per_char))
    lines = 0
    for para in text.split('\n'):
        lines += max(1, -(-len(para) // cols))
    return lines


for idx, slide in enumerate(prs.slides, start=1):
    boxes = []
    for sh in slide.shapes:
        if sh.left is None:
            continue
        x, y, w, h = box(sh)

        # 1. Escaping the slide.
        if x < -0.01 or y < -0.01 or x + w > SLIDE_W + 0.01 or y + h > SLIDE_H + 0.01:
            issues.append(
                f'slide {idx}: shape escapes the slide '
                f'({x:.2f},{y:.2f} {w:.2f}x{h:.2f})'
            )

        # 2. Text that will not fit its box.
        if sh.has_text_frame and sh.text_frame.text.strip():
            tf = sh.text_frame
            text = tf.text
            size = 12.0
            face = 'Calibri'
            for p in tf.paragraphs:
                for r in p.runs:
                    if r.font.size:
                        size = r.font.size.pt
                    if r.font.name:
                        face = r.font.name
                    break
                break
            # Use the paragraph's own line spacing when set, else 1.22x.
            spacing = None
            for p in tf.paragraphs:
                if p.line_spacing:
                    spacing = (p.line_spacing.pt / 72.0
                               if hasattr(p.line_spacing, 'pt')
                               else p.line_spacing * size / 72.0)
                break
            line_h = spacing or (size * 1.22 / 72.0)
            need = est_lines(text, max(w - 0.12, 0.3), size, face) * line_h
            if need > h + 0.10:
                issues.append(
                    f'slide {idx}: text may overflow — needs ~{need:.2f}" in {h:.2f}" '
                    f'@{size:.0f}pt :: {text.strip()[:58]!r}'
                )
            boxes.append((x, y, w, h, text.strip()[:34]))

    # 3. Card-on-card overlap (text boxes are allowed to sit on cards).
    fills = []
    for sh in slide.shapes:
        if sh.left is None or sh.has_text_frame and sh.text_frame.text.strip():
            continue
        x, y, w, h = box(sh)
        if w > 0.9 and h > 0.5:      # only real panels, not rules or dots
            fills.append((x, y, w, h))
    for i in range(len(fills)):
        for j in range(i + 1, len(fills)):
            ax, ay, aw, ah = fills[i]
            bx, by, bw, bh = fills[j]
            ox = min(ax + aw, bx + bw) - max(ax, bx)
            oy = min(ay + ah, by + bh) - max(ay, by)
            if ox > 0.06 and oy > 0.06:
                issues.append(
                    f'slide {idx}: panels overlap by {ox:.2f}x{oy:.2f}" '
                    f'({ax:.2f},{ay:.2f}) vs ({bx:.2f},{by:.2f})'
                )

print(f'{len(prs.slides)} slides checked')
if issues:
    print(f'\n{len(issues)} issue(s):')
    for i in issues:
        print('  -', i)
    sys.exit(1)
print('no geometry issues found')
